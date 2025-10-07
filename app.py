from flask import Flask, request, render_template, jsonify
import os, re
from docx import Document
import pdfplumber
from pptx import Presentation
from openpyxl import load_workbook
import traceback

app = Flask(__name__)

UPLOAD_FOLDER = "uploads"
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

# Add this for Render deployment
PORT = int(os.environ.get('PORT', 10000))

# Set memory limits for PDF processing
PDF_MEMORY_LIMIT = 50 * 1024 * 1024  # 50MB

def search_in_docx(file_path, query):
    results = []
    try:
        doc = Document(file_path)
        for para in doc.paragraphs:
            if query.lower() in para.text.lower():
                results.append(para.text.strip())
    except Exception as e:
        print(f"Error reading {file_path}: {e}")
    return results

def search_in_pdf(file_path, query):
    results = []
    try:
        # Check file size before processing (limit to 10MB for free tier)
        file_size = os.path.getsize(file_path)
        if file_size > 10 * 1024 * 1024:  # 10MB limit
            return [f"PDF too large to process ({file_size//1024//1024}MB). Please use smaller files."]
        
        with pdfplumber.open(file_path) as pdf:
            for page_num, page in enumerate(pdf.pages[:5], start=1):  # Limit to first 5 pages
                try:
                    text = page.extract_text()
                    if text and query.lower() in text.lower():
                        match = re.search(f".{{0,30}}{re.escape(query)}.{{0,30}}", text, re.IGNORECASE)
                        if match:
                            results.append(f"[Page {page_num}] ...{match.group(0)}...")
                except Exception as e:
                    print(f"Error processing page {page_num}: {e}")
                    continue
    except Exception as e:
        print(f"Error reading {file_path}: {e}")
        results.append(f"Error processing PDF: {str(e)}")
    return results

def search_in_pptx(file_path, query):
    results = []
    try:
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides[:10], start=1):  # Limit to first 10 slides
            for shape in slide.shapes:
                if hasattr(shape, "text") and query.lower() in shape.text.lower():
                    snippet = re.search(f".{{0,30}}{re.escape(query)}.{{0,30}}", shape.text, re.IGNORECASE)
                    if snippet:
                        results.append(f"[Slide {i}] ...{snippet.group(0)}...")
    except Exception as e:
        print(f"Error reading {file_path}: {e}")
    return results

def search_in_excel(file_path, query):
    results = []
    try:
        wb = load_workbook(file_path, data_only=True, read_only=True)  # Use read_only mode
        for sheet in wb.sheetnames[:3]:  # Limit to first 3 sheets
            ws = wb[sheet]
            for row in ws.iter_rows(max_row=1000, values_only=True):  # Limit rows
                for cell in row:
                    if cell and isinstance(cell, str) and query.lower() in cell.lower():
                        results.append(f"[{sheet}] {cell.strip()}")
        wb.close()
    except Exception as e:
        print(f"Error reading {file_path}: {e}")
    return results

@app.route('/')
def home():
    return render_template('index.html')

@app.route('/search', methods=['POST'])
def search():
    query = request.form.get('query', '').strip()
    if not query:
        return render_template("index.html", error="Please enter a search query")
    
    results = []
    
    try:
        for root, _, files in os.walk(UPLOAD_FOLDER):
            for file in files:
                file_path = os.path.join(root, file)
                ext = os.path.splitext(file)[1].lower()
                file_results = []

                try:
                    if ext == ".docx":
                        file_results = search_in_docx(file_path, query)
                    elif ext == ".pdf":
                        file_results = search_in_pdf(file_path, query)
                    elif ext == ".pptx":
                        file_results = search_in_pptx(file_path, query)
                    elif ext in [".xlsx", ".xls"]:
                        file_results = search_in_excel(file_path, query)

                    if file_results:
                        results.append({"file": file, "matches": file_results})
                except Exception as e:
                    print(f"Error processing {file}: {e}")
                    continue

    except Exception as e:
        print(f"Search error: {e}")
        return render_template("index.html", error=f"Search error: {str(e)}")

    return render_template("index.html", query=query, results=results)

@app.errorhandler(500)
def internal_error(error):
    return render_template("index.html", error="Internal server error. The file might be too large or corrupted.")

@app.errorhandler(404)
def not_found(error):
    return render_template("index.html", error="Page not found.")

if __name__ == "__main__":
    app.run(host='0.0.0.0', port=PORT, debug=False)
